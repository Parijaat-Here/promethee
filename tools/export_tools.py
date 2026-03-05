# ============================================================================
# Prométhée — Assistant IA desktop
# ============================================================================
# Auteur  : Pierre COUGET
# Licence : GNU Affero General Public License v3.0 (AGPL-3.0)
#           https://www.gnu.org/licenses/agpl-3.0.html
# Année   : 2026
# ----------------------------------------------------------------------------
# Ce fichier fait partie du projet Prométhée.
# Vous pouvez le redistribuer et/ou le modifier selon les termes de la
# licence AGPL-3.0 publiée par la Free Software Foundation.
# ============================================================================

"""
tools/export_tools.py — Génération et export de fichiers bureautiques
======================================================================

Outils exposés (9) :

  Markdown (1) :
    - export_md         : écrit un fichier Markdown à partir de contenu texte

  Word / docx (1) :
    - export_docx       : génère un document Word structuré (titres, paragraphes,
                          tableaux, listes) depuis une description JSON

  Tableur / xlsx (2) :
    - export_xlsx_json  : génère un classeur Excel depuis une structure JSON
                          (feuilles, en-têtes, lignes, graphiques)
    - export_xlsx_csv   : génère un classeur Excel depuis du CSV brut (une feuille)

  Présentation / pptx (2) :
    - export_pptx_json  : génère une présentation PowerPoint depuis une structure
                          JSON (titre, slides, puces, notes)
    - export_pptx_outline: génère une présentation depuis un outline texte
                          (ligne "# Titre", "- Puce", "> Note")

  PDF (1) :
    - export_pdf        : génère un PDF structuré (titres, paragraphes, tableaux)
                          depuis une description JSON, via reportlab

  LibreOffice natif (2) :
    - export_libreoffice: convertit un fichier existant vers odt/ods/odp
                          en invoquant LibreOffice headless
    - export_libreoffice_native : génère directement un odt/ods/odp depuis
                          une description JSON (via python-docx/openpyxl/python-pptx
                          + conversion LibreOffice)

Conventions communes
────────────────────
  - output_path : chemin absolu ou relatif au home utilisateur.
                  Si omis ou vide, un fichier est créé dans ~/Exports/Prométhée/.
  - Retour      : dict JSON {"path": "/chemin/absolu", "size_bytes": N,
                             "pages"/"sheets"/"slides": N, "status": "ok"}
  - En cas d'erreur : {"error": "message explicatif", "status": "error"}

Structure JSON commune pour export_docx, export_pdf, export_pptx_json
──────────────────────────────────────────────────────────────────────
  {
    "title": "Titre du document",
    "sections": [
      {
        "heading": "Titre de section",   // niveau 1-3 (optionnel)
        "level": 1,                      // 1, 2 ou 3 (défaut 1)

        // Contenu textuel — trois variantes (choisir la plus adaptée) :
        "paragraphs": ["Para 1 développé.", "Para 2 développé."],
                                         // PRÉFÉRÉ pour plusieurs paragraphes
        "content": "Texte avec\n\ndoubles sauts\n\nou simple paragraphe",
                                         // alternatif à paragraphs
        "intro": "Texte avant bullets/tableau",  // cumulable avec bullets/table

        // Éléments structurés (cumulables avec intro) :
        "table": {
          "headers": ["Col A", "Col B"],
          "rows": [["val1", "val2"], ["val3", "val4"]]
        },
        "bullets": ["item 1", "item 2"],

        "page_break": false              // saut de page optionnel après la section
      }
    ]
  }

  NOTE : un document professionnel complet comporte typiquement 15 à 40 sections.
  Chaque section doit contenir un contenu rédigé, dense et proportionnel au sujet.

Prérequis :
    pip install python-docx openpyxl python-pptx reportlab
    LibreOffice installé système (apt install libreoffice)
"""

import io
import json
import os
import re
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any

from core.tools_engine import tool, set_current_family

set_current_family("export_tools", "Export de fichiers", "📄")

# ── Répertoire de sortie par défaut ──────────────────────────────────────────

_DEFAULT_EXPORT_DIR = Path.home() / "Exports" / "Prométhée"


def _resolve_output(output_path: str, default_name: str) -> Path:
    """Résout le chemin de sortie et crée les répertoires manquants."""
    if output_path:
        p = Path(output_path).expanduser()
        if not p.is_absolute():
            p = Path.home() / p
    else:
        _DEFAULT_EXPORT_DIR.mkdir(parents=True, exist_ok=True)
        p = _DEFAULT_EXPORT_DIR / default_name
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


def _ok(path: Path, extra: dict | None = None) -> str:
    r = {"status": "ok", "path": str(path), "size_bytes": path.stat().st_size}
    if extra:
        r.update(extra)
    return json.dumps(r, ensure_ascii=False)


def _err(msg: str) -> str:
    return json.dumps({"status": "error", "error": msg}, ensure_ascii=False)


# ═════════════════════════════════════════════════════════════════════════════
# MARKDOWN
# ═════════════════════════════════════════════════════════════════════════════

@tool(
    name="export_md",
    description=(
        "Écrit un fichier Markdown (.md) à partir de contenu texte brut. "
        "À utiliser quand l'utilisateur demande de produire un rapport, un README, "
        "un article ou tout document texte en format Markdown. "
        "Le contenu doit déjà être en Markdown valide (titres #, listes -, tableaux |)."
    ),
    parameters={
        "type": "object",
        "properties": {
            "content": {
                "type": "string",
                "description": "Contenu Markdown complet à écrire dans le fichier."
            },
            "output_path": {
                "type": "string",
                "description": (
                    "Chemin de destination, ex: ~/Documents/rapport.md. "
                    "Si omis, crée le fichier dans ~/Exports/Prométhée/."
                )
            },
            "filename": {
                "type": "string",
                "description": "Nom du fichier si output_path est omis (ex: rapport.md)."
            }
        },
        "required": ["content"]
    }
)
def export_md(content: str, output_path: str = "", filename: str = "") -> str:
    try:
        name = filename or "export.md"
        if not name.endswith(".md"):
            name += ".md"
        p = _resolve_output(output_path, name)
        p.write_text(content, encoding="utf-8")
        lines = content.count("\n") + 1
        return _ok(p, {"lines": lines})
    except Exception as e:
        return _err(f"export_md : {e}")


# ═════════════════════════════════════════════════════════════════════════════
# WORD / DOCX
# ═════════════════════════════════════════════════════════════════════════════

def _build_docx(doc_json: dict):
    """Construit un docx à partir de la structure JSON commune.

    Améliorations v2 :
    - 'paragraphs' : liste de chaînes → plusieurs paragraphes dans une section
    - 'content' accepte les sauts de ligne (\n) comme séparateurs de paragraphes
    - 'intro' : paragraphe introductif avant les bullets ou le tableau
    - 'bullets' et 'table' peuvent coexister avec 'content'/'intro'
    - Compatibilité totale avec les documents générés avant cette version
    """
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Titre principal
    if doc_json.get("title"):
        doc.add_heading(doc_json["title"], level=0)

    page_count = 1
    for section in doc_json.get("sections", []):
        heading = section.get("heading")
        level   = max(1, min(3, int(section.get("level", 1))))

        if heading:
            doc.add_heading(heading, level=level)

        # ── Paragraphes ────────────────────────────────────────────────
        # 1. 'paragraphs' : liste explicite de paragraphes (prioritaire)
        if section.get("paragraphs"):
            for para in section["paragraphs"]:
                if para and str(para).strip():
                    doc.add_paragraph(str(para))

        # 2. 'content' : chaîne simple ou multi-lignes (séparées par \n\n ou \n)
        elif section.get("content"):
            raw = section["content"]
            # Découper sur les doubles sauts de ligne (paragraphes distincts)
            # puis sur les sauts simples si aucun double n'est présent
            if "\n\n" in raw:
                parts = [p.strip() for p in raw.split("\n\n") if p.strip()]
            else:
                parts = [p.strip() for p in raw.split("\n") if p.strip()]
            for part in parts:
                doc.add_paragraph(part)

        # 3. 'intro' : paragraphe placé avant les bullets ou le tableau
        #    (compatibilité avec les sections qui ont intro + bullets)
        if section.get("intro"):
            doc.add_paragraph(str(section["intro"]))

        # ── Liste à puces ───────────────────────────────────────────────
        if section.get("bullets"):
            for item in section["bullets"]:
                if item and str(item).strip():
                    doc.add_paragraph(str(item), style="List Bullet")

        # ── Tableau ─────────────────────────────────────────────────────
        if section.get("table"):
            tbl_data = section["table"]
            headers = tbl_data.get("headers", [])
            rows    = tbl_data.get("rows", [])
            if headers:
                table = doc.add_table(rows=1 + len(rows), cols=len(headers))
                table.style = "Light Grid Accent 1"
                hdr_cells = table.rows[0].cells
                for i, h in enumerate(headers):
                    hdr_cells[i].text = str(h)
                for ri, row in enumerate(rows):
                    row_cells = table.rows[ri + 1].cells
                    for ci, val in enumerate(row[:len(headers)]):
                        row_cells[ci].text = str(val)
                doc.add_paragraph()  # espace après tableau

        # Page break optionnel
        if section.get("page_break"):
            doc.add_page_break()
            page_count += 1

    return doc


@tool(
    name="export_docx",
    description=(
        "Génère un document Word (.docx) structuré depuis une description JSON. "
        "Supporte les titres hiérarchiques (niveaux 1 à 3), les listes à puces, les tableaux "
        "et plusieurs paragraphes par section. "
        "À utiliser pour tout document formel : rapport, compte-rendu, contrat, note, guide. "
        "\n\n"
        "RÈGLE DE CONTENU CRITIQUE : le document doit être COMPLET et DÉVELOPPÉ. "
        "Ne pas se limiter à une structure squelette. "
        "Chaque section doit contenir un vrai contenu rédigé, proportionnel au sujet. "
        "Un rapport sérieux comporte typiquement 15 à 40 sections. "
        "Utiliser 'paragraphs' (liste) pour plusieurs paragraphes dans une section, "
        "ou 'content' avec des doubles sauts de ligne (\\n\\n) comme séparateurs. "
        "'bullets' et 'table' peuvent être combinés avec 'intro' dans la même section. "
        "\n\n"
        "Structure JSON d'une section complète : "
        '{"heading": "Titre", "level": 1, '
        '"intro": "Texte introductif (optionnel, avant les bullets/tableau)", '
        '"paragraphs": ["Paragraphe 1 développé.", "Paragraphe 2 développé."], '
        '"bullets": ["Point clé 1", "Point clé 2"], '
        '"table": {"headers": ["Col A", "Col B"], "rows": [["v1","v2"]]}, '
        '"page_break": false}'
    ),
    parameters={
        "type": "object",
        "properties": {
            "document": {
                "type": "object",
                "description": (
                    "Structure du document Word à générer. "
                    "Champ 'title' (str, titre principal). "
                    "Champ 'sections' : liste d'objets section. "
                    "Chaque section accepte : "
                    "heading (str, titre de section), "
                    "level (int 1-3, niveau de titre, défaut 1), "
                    "paragraphs (liste de str — PRÉFÉRER à 'content' pour plusieurs paragraphes), "
                    "content (str — paragraphe unique ou texte avec \\n\\n pour multi-paragraphes), "
                    "intro (str — paragraphe placé avant bullets ou tableau), "
                    "bullets (liste de str — puces, cumulable avec intro), "
                    "table (objet {headers, rows} — tableau, cumulable avec intro), "
                    "page_break (bool — saut de page après la section, défaut false). "
                    "IMPORTANT : rédiger un contenu dense et complet ; "
                    "un document de qualité professionnelle comporte 15 à 40 sections."
                )
            },
            "output_path": {
                "type": "string",
                "description": "Chemin de destination (ex: ~/Documents/rapport.docx). Optionnel."
            },
            "filename": {
                "type": "string",
                "description": "Nom du fichier si output_path est omis."
            }
        },
        "required": ["document"]
    }
)
def export_docx(document: dict, output_path: str = "", filename: str = "") -> str:
    try:
        name = filename or (document.get("title", "export") + ".docx")
        if not name.endswith(".docx"):
            name += ".docx"
        p = _resolve_output(output_path, name)
        doc = _build_docx(document)
        doc.save(str(p))
        sections = len(document.get("sections", []))
        return _ok(p, {"sections": sections})
    except Exception as e:
        return _err(f"export_docx : {e}")


# ═════════════════════════════════════════════════════════════════════════════
# EXCEL / XLSX
# ═════════════════════════════════════════════════════════════════════════════

# ── Types de graphiques supportés ────────────────────────────────────────────
#
#  bar        Barres verticales groupées       (comparaison entre catégories)
#  bar_stacked Barres verticales empilées      (composition + total)
#  bar_percent Barres empilées 100 %           (proportions)
#  bar_h       Barres horizontales groupées    (libellés longs)
#  line        Courbes                         (évolution temporelle)
#  line_smooth Courbes lissées (spline)        (tendances)
#  area        Aires                           (volumes cumulés)
#  area_stacked Aires empilées                 (parts cumulées)
#  pie         Camembert                       (parts d'un tout, ≤ 6 séries)
#  doughnut    Anneau                          (idem camembert, style moderne)
#  scatter     Nuage de points                 (corrélations entre 2 variables)
#  bubble      Bulles                          (3 variables : x, y, taille)
#  radar       Radar / toile d'araignée        (profils multi-critères)
#
# ── Structure d'un graphique dans sheets[].charts[] ─────────────────────────
#
#  {
#    "type": "bar",                  # obligatoire — voir liste ci-dessus
#    "title": "Ventes par mois",     # titre affiché au-dessus du graphique
#    "data_sheet": "Ventes",         # feuille source (défaut : feuille courante)
#    "categories_col": 1,            # colonne des étiquettes X (indice 1-based)
#    "series": [                     # liste des séries à tracer
#      {
#        "title": "CA (€)",          # légende de la série
#        "col": 2,                   # colonne des valeurs (indice 1-based)
#        "color": "2255A4"           # couleur hex optionnelle (sans #)
#      }
#    ],
#    "data_rows": [2, 13],           # [première_ligne, dernière_ligne] données
#                                    # si omis : toutes les lignes après l'en-tête
#    "anchor": "E2",                 # cellule d'ancrage du coin supérieur gauche
#                                    # (défaut : colonne après les données, ligne 1)
#    "width_cm": 15,                 # largeur en cm (défaut : 15)
#    "height_cm": 10,                # hauteur en cm (défaut : 10)
#    "style": 10,                    # style Excel 1-48 (défaut : 10)
#    "show_legend": true,            # afficher la légende (défaut : true)
#    "show_data_labels": false       # afficher les valeurs sur les barres (défaut : false)
#  }

def _build_charts(wb, ws, charts_def: list, sheet_data_rows: int, sheet_headers_count: int) -> list[str]:
    """
    Construit les graphiques openpyxl pour une feuille.
    Retourne la liste des avertissements éventuels.
    """
    from openpyxl.chart import (
        BarChart, LineChart, AreaChart, PieChart, DoughnutChart,
        ScatterChart, BubbleChart, RadarChart,
        Reference, Series,
    )
    from openpyxl.chart.series import SeriesLabel
    from openpyxl.chart.label import DataLabel
    from openpyxl.utils import get_column_letter
    import openpyxl.utils.units as units_util

    warnings = []

    _CHART_FACTORIES = {
        "bar":          lambda: _make_bar(grouping="clustered",   direction="col"),
        "bar_stacked":  lambda: _make_bar(grouping="stacked",     direction="col"),
        "bar_percent":  lambda: _make_bar(grouping="percentStacked", direction="col"),
        "bar_h":        lambda: _make_bar(grouping="clustered",   direction="bar"),
        "line":         lambda: LineChart(),
        "line_smooth":  lambda: _make_line_smooth(),
        "area":         lambda: AreaChart(),
        "area_stacked": lambda: _make_area_stacked(),
        "pie":          lambda: PieChart(),
        "doughnut":     lambda: DoughnutChart(),
        "scatter":      lambda: ScatterChart(),
        "bubble":       lambda: BubbleChart(),
        "radar":        lambda: RadarChart(),
    }

    def _make_bar(grouping: str, direction: str):
        c = BarChart()
        c.type = direction
        c.grouping = grouping
        return c

    def _make_line_smooth():
        c = LineChart()
        c.smooth = True
        return c

    def _make_area_stacked():
        c = AreaChart()
        c.grouping = "stacked"
        return c

    for idx, chart_def in enumerate(charts_def):
        chart_type = str(chart_def.get("type", "bar")).lower()
        if chart_type not in _CHART_FACTORIES:
            warnings.append(f"Graphique {idx+1} : type inconnu '{chart_type}', ignoré.")
            continue

        try:
            chart = _CHART_FACTORIES[chart_type]()

            # Titre
            if chart_def.get("title"):
                chart.title = chart_def["title"]

            # Style
            style = chart_def.get("style", 10)
            if isinstance(style, int) and 1 <= style <= 48:
                chart.style = style

            # Légende
            if not chart_def.get("show_legend", True):
                chart.legend = None

            # Dimensions
            w_cm = chart_def.get("width_cm", 15)
            h_cm = chart_def.get("height_cm", 10)
            chart.width  = w_cm
            chart.height = h_cm

            # Feuille source des données
            src_sheet_name = chart_def.get("data_sheet")
            src_ws = wb[src_sheet_name] if src_sheet_name and src_sheet_name in wb.sheetnames else ws

            # Plage de lignes des données
            data_rows = chart_def.get("data_rows")
            if data_rows and len(data_rows) == 2:
                row_min, row_max = int(data_rows[0]), int(data_rows[1])
            else:
                row_min, row_max = 2, sheet_data_rows + 1

            # Catégories (axe X / libellés)
            cat_col = chart_def.get("categories_col", 1)
            cats = Reference(
                src_ws,
                min_col=cat_col, max_col=cat_col,
                min_row=row_min, max_row=row_max,
            )

            # Séries
            series_defs = chart_def.get("series", [])
            if not series_defs:
                warnings.append(f"Graphique {idx+1} ('{chart_def.get('title', '')}') : aucune série définie, ignoré.")
                continue

            for s_def in series_defs:
                col = int(s_def.get("col", 2))
                vals = Reference(
                    src_ws,
                    min_col=col, max_col=col,
                    min_row=row_min, max_row=row_max,
                )

                if chart_type in ("scatter", "bubble"):
                    # ScatterChart / BubbleChart : x_values + values
                    x_vals = Reference(src_ws, min_col=cat_col, max_col=cat_col,
                                       min_row=row_min, max_row=row_max)
                    series = Series(vals, xvalues=x_vals)
                else:
                    series = Series(vals, cats)

                # Titre de la série (légende)
                if s_def.get("title"):
                    series.title = SeriesLabel(v=s_def["title"])

                # Couleur de remplissage
                if s_def.get("color"):
                    from openpyxl.drawing.fill import PatternFillProperties
                    from openpyxl.drawing.spreadsheet_drawing import SpreadsheetDrawing
                    try:
                        hex_color = s_def["color"].lstrip("#")
                        series.graphicalProperties.solidFill = hex_color
                    except Exception:
                        pass

                chart.series.append(series)

            # Étiquettes de données
            if chart_def.get("show_data_labels", False):
                try:
                    chart.dLbls = DataLabel()
                    chart.dLbls.showVal = True
                    chart.dLbls.showLegendKey = False
                    chart.dLbls.showCatName = False
                    chart.dLbls.showSerName = False
                except Exception:
                    pass

            # Cellule d'ancrage
            anchor = chart_def.get("anchor")
            if not anchor:
                # Par défaut : première colonne après les données, ligne 1
                next_col = (max((int(s.get("col", 2)) for s in series_defs), default=2) + 2)
                anchor = f"{get_column_letter(next_col)}1"

            ws.add_chart(chart, anchor)

        except Exception as e:
            warnings.append(f"Graphique {idx+1} ('{chart_def.get('title', '')}') : erreur — {e}")

    return warnings


@tool(
    name="export_xlsx_json",
    description=(
        "Génère un classeur Excel (.xlsx) depuis une structure JSON. "
        "Permet de créer plusieurs feuilles avec en-têtes, données et graphiques. "
        "Chaque feuille peut contenir un champ optionnel 'charts' pour générer "
        "des graphiques directement intégrés dans le classeur. "
        "Types de graphiques : bar, bar_stacked, bar_percent, bar_h, line, line_smooth, "
        "area, area_stacked, pie, doughnut, scatter, bubble, radar."
    ),
    parameters={
        "type": "object",
        "properties": {
            "workbook": {
                "type": "object",
                "description": (
                    "Structure du classeur. Champ 'sheets' : liste d'objets avec "
                    "name (str), headers (liste de str), rows (liste de listes), "
                    "et charts (liste optionnelle de graphiques). "
                    "Chaque graphique : {type, title, categories_col, series: [{title, col, color}], "
                    "data_rows, anchor, width_cm, height_cm, style, show_legend, show_data_labels}. "
                    "Exemple : voir la skill guide_export_excel.md pour la structure complète."
                )
            },
            "output_path": {
                "type": "string",
                "description": "Chemin de destination (ex: ~/Documents/données.xlsx). Optionnel."
            },
            "filename": {
                "type": "string",
                "description": "Nom du fichier si output_path est omis."
            }
        },
        "required": ["workbook"]
    }
)
def export_xlsx_json(workbook: dict, output_path: str = "", filename: str = "") -> str:
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment
        from openpyxl.utils import get_column_letter

        sheets = workbook.get("sheets", [])
        if not sheets:
            return _err("export_xlsx_json : 'sheets' est vide ou manquant")

        name = filename or "export.xlsx"
        if not name.endswith(".xlsx"):
            name += ".xlsx"
        p = _resolve_output(output_path, name)

        wb = openpyxl.Workbook()
        wb.remove(wb.active)  # supprimer la feuille vide par défaut

        header_font    = Font(bold=True, color="FFFFFF")
        header_fill    = PatternFill("solid", fgColor="2255A4")
        header_align   = Alignment(horizontal="center", vertical="center")
        alt_fill       = PatternFill("solid", fgColor="EFF3FA")

        total_rows  = 0
        all_warnings = []

        for sheet_def in sheets:
            ws = wb.create_sheet(title=str(sheet_def.get("name", "Feuille"))[:31])
            headers = sheet_def.get("headers", [])
            rows    = sheet_def.get("rows", [])

            # En-têtes
            for ci, h in enumerate(headers, 1):
                cell = ws.cell(row=1, column=ci, value=str(h))
                cell.font      = header_font
                cell.fill      = header_fill
                cell.alignment = header_align

            # Données + formatage alterné
            for ri, row in enumerate(rows, 2):
                fill = alt_fill if ri % 2 == 0 else None
                for ci, val in enumerate(row[:len(headers) or len(row)], 1):
                    cell = ws.cell(row=ri, column=ci, value=val)
                    if fill:
                        cell.fill = fill

            # Ajustement auto de la largeur des colonnes
            for col_cells in ws.columns:
                max_len = max((len(str(c.value or "")) for c in col_cells), default=8)
                ws.column_dimensions[get_column_letter(col_cells[0].column)].width = min(max_len + 4, 60)

            # Figer la première ligne
            ws.freeze_panes = "A2"
            total_rows += len(rows)

            # Graphiques
            charts_def = sheet_def.get("charts", [])
            if charts_def:
                warnings = _build_charts(wb, ws, charts_def, len(rows), len(headers))
                all_warnings.extend(warnings)

        wb.save(str(p))
        result = {"sheets": len(sheets), "total_rows": total_rows}
        if all_warnings:
            result["warnings"] = all_warnings
        return _ok(p, result)
    except Exception as e:
        return _err(f"export_xlsx_json : {e}")


@tool(
    name="export_xlsx_csv",
    description=(
        "Génère un fichier Excel (.xlsx) depuis du contenu CSV brut. "
        "Pratique quand les données sont déjà disponibles en format CSV "
        "(séparateur virgule ou point-virgule détecté automatiquement). "
        "Crée une seule feuille. Pour plusieurs feuilles, utiliser export_xlsx_json."
    ),
    parameters={
        "type": "object",
        "properties": {
            "csv_content": {
                "type": "string",
                "description": "Contenu CSV brut (avec en-têtes sur la première ligne)."
            },
            "sheet_name": {
                "type": "string",
                "description": "Nom de la feuille (défaut : 'Données')."
            },
            "output_path": {
                "type": "string",
                "description": "Chemin de destination. Optionnel."
            },
            "filename": {
                "type": "string",
                "description": "Nom du fichier si output_path est omis."
            }
        },
        "required": ["csv_content"]
    }
)
def export_xlsx_csv(csv_content: str, sheet_name: str = "Données",
                    output_path: str = "", filename: str = "") -> str:
    try:
        import csv as csv_mod
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment
        from openpyxl.utils import get_column_letter

        # Détecter le séparateur
        sample = csv_content[:2000]
        dialect = csv_mod.Sniffer().sniff(sample, delimiters=",;\t|")
        reader  = csv_mod.reader(io.StringIO(csv_content), dialect)
        all_rows = list(reader)

        if not all_rows:
            return _err("export_xlsx_csv : CSV vide")

        name = filename or "export.xlsx"
        if not name.endswith(".xlsx"):
            name += ".xlsx"
        p = _resolve_output(output_path, name)

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = str(sheet_name)[:31]

        header_font  = Font(bold=True, color="FFFFFF")
        header_fill  = PatternFill("solid", fgColor="2255A4")
        header_align = Alignment(horizontal="center")
        alt_fill     = PatternFill("solid", fgColor="EFF3FA")

        for ri, row in enumerate(all_rows, 1):
            for ci, val in enumerate(row, 1):
                # Tenter une conversion numérique
                cell = ws.cell(row=ri, column=ci)
                try:
                    cell.value = int(val)
                except ValueError:
                    try:
                        cell.value = float(val.replace(",", "."))
                    except ValueError:
                        cell.value = val

                if ri == 1:
                    cell.font      = header_font
                    cell.fill      = header_fill
                    cell.alignment = header_align
                elif ri % 2 == 0:
                    cell.fill = alt_fill

        # Ajustement largeur
        for col_cells in ws.columns:
            max_len = max((len(str(c.value or "")) for c in col_cells), default=8)
            ws.column_dimensions[get_column_letter(col_cells[0].column)].width = min(max_len + 4, 60)

        ws.freeze_panes = "A2"
        wb.save(str(p))
        return _ok(p, {"rows": len(all_rows) - 1, "columns": len(all_rows[0]) if all_rows else 0})
    except Exception as e:
        return _err(f"export_xlsx_csv : {e}")


# ═════════════════════════════════════════════════════════════════════════════
# POWERPOINT / PPTX
# ═════════════════════════════════════════════════════════════════════════════

def _add_slide(prs, layout_idx: int, title: str, bullets: list[str],
               content: str = "", notes: str = "", subtitle: str = ""):
    """Ajoute un slide à la présentation."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    layout = prs.slide_layouts[layout_idx]
    slide  = prs.slides.add_slide(layout)

    # Titre
    if slide.shapes.title:
        slide.shapes.title.text = title

    # Sous-titre (layout 0 = titre principal)
    if layout_idx == 0 and subtitle:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                ph.text = subtitle
                break

    # Corps / puces
    if bullets or content:
        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            if idx in (1, 2):  # corps du slide
                tf = ph.text_frame
                tf.clear()
                if bullets:
                    for i, b in enumerate(bullets):
                        if i == 0:
                            tf.paragraphs[0].text = b
                            tf.paragraphs[0].level = 0
                        else:
                            p = tf.add_paragraph()
                            p.text  = b
                            p.level = 0
                elif content:
                    tf.paragraphs[0].text = content
                break

    # Notes
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


@tool(
    name="export_pptx_json",
    description=(
        "Génère une présentation PowerPoint (.pptx) depuis une structure JSON. "
        "Chaque slide peut avoir un titre, une liste de puces, du contenu texte "
        "et des notes de présentation. "
        "À utiliser pour créer des présentations, des pitch decks, des diaporamas."
    ),
    parameters={
        "type": "object",
        "properties": {
            "presentation": {
                "type": "object",
                "description": (
                    "Structure de la présentation. Champs : "
                    "title (str, titre global), "
                    "subtitle (str, optionnel, pour le slide de titre), "
                    "slides (liste d'objets avec title/bullets/content/notes). "
                    "Exemple : {\"title\": \"Mon rapport\", \"slides\": ["
                    "{\"title\": \"Introduction\", \"bullets\": [\"Point 1\", \"Point 2\"], "
                    "\"notes\": \"Penser à mentionner le contexte.\"}]}"
                )
            },
            "output_path": {
                "type": "string",
                "description": "Chemin de destination (ex: ~/Documents/présentation.pptx). Optionnel."
            },
            "filename": {
                "type": "string",
                "description": "Nom du fichier si output_path est omis."
            }
        },
        "required": ["presentation"]
    }
)
def export_pptx_json(presentation: dict, output_path: str = "", filename: str = "") -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        name = filename or (presentation.get("title", "présentation") + ".pptx")
        if not name.endswith(".pptx"):
            name += ".pptx"
        p = _resolve_output(output_path, name)

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Slide de titre
        main_title = presentation.get("title", "")
        subtitle   = presentation.get("subtitle", "")
        if main_title:
            _add_slide(prs, 0, main_title, [], subtitle=subtitle)

        # Slides de contenu
        for slide_def in presentation.get("slides", []):
            title   = slide_def.get("title", "")
            bullets = slide_def.get("bullets", [])
            content = slide_def.get("content", "")
            notes   = slide_def.get("notes", "")
            _add_slide(prs, 1, title, bullets, content=content, notes=notes)

        prs.save(str(p))
        n_slides = len(prs.slides)
        return _ok(p, {"slides": n_slides})
    except Exception as e:
        return _err(f"export_pptx_json : {e}")


@tool(
    name="export_pptx_outline",
    description=(
        "Génère une présentation PowerPoint (.pptx) depuis un outline texte simple. "
        "Format : '# Titre du slide' pour un nouveau slide, '- Puce' pour une puce, "
        "'> Note' pour une note de présentation, texte libre pour contenu. "
        "Plus rapide à écrire que le format JSON pour des présentations simples."
    ),
    parameters={
        "type": "object",
        "properties": {
            "outline": {
                "type": "string",
                "description": (
                    "Outline texte de la présentation. "
                    "Exemple :\\n"
                    "# Introduction\\n"
                    "- Contexte du projet\\n"
                    "- Objectifs\\n"
                    "> Penser à présenter l'équipe\\n"
                    "# Résultats\\n"
                    "- Hausse de 15% du CA\\n"
                )
            },
            "title": {
                "type": "string",
                "description": "Titre global de la présentation (slide de couverture)."
            },
            "output_path": {
                "type": "string",
                "description": "Chemin de destination. Optionnel."
            },
            "filename": {
                "type": "string",
                "description": "Nom du fichier si output_path est omis."
            }
        },
        "required": ["outline"]
    }
)
def export_pptx_outline(outline: str, title: str = "",
                         output_path: str = "", filename: str = "") -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches

        name = filename or (title or "présentation") + ".pptx"
        if not name.endswith(".pptx"):
            name += ".pptx"
        p = _resolve_output(output_path, name)

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Slide de couverture
        if title:
            _add_slide(prs, 0, title, [])

        # Parser l'outline
        slides: list[dict] = []
        current: dict | None = None

        for raw_line in outline.splitlines():
            line = raw_line.rstrip()
            if line.startswith("# "):
                if current is not None:
                    slides.append(current)
                current = {"title": line[2:].strip(), "bullets": [], "notes": "", "content": ""}
            elif line.startswith("- ") and current is not None:
                current["bullets"].append(line[2:].strip())
            elif line.startswith("> ") and current is not None:
                current["notes"] += line[2:].strip() + " "
            elif line.strip() and current is not None and not current["bullets"]:
                current["content"] += line.strip() + " "

        if current is not None:
            slides.append(current)

        for s in slides:
            _add_slide(prs, 1, s["title"], s["bullets"],
                       content=s["content"].strip(), notes=s["notes"].strip())

        prs.save(str(p))
        return _ok(p, {"slides": len(prs.slides)})
    except Exception as e:
        return _err(f"export_pptx_outline : {e}")


# ═════════════════════════════════════════════════════════════════════════════
# PDF
# ═════════════════════════════════════════════════════════════════════════════

@tool(
    name="export_pdf",
    description=(
        "Génère un document PDF structuré depuis une description JSON. "
        "Supporte les titres hiérarchiques (niveaux 1 à 3), les paragraphes, "
        "les listes à puces et les tableaux. "
        "À utiliser pour produire des rapports, des fiches, des documents formels en PDF. "
        "Même format JSON que export_docx : utiliser 'paragraphs' (liste) pour plusieurs "
        "paragraphes, 'intro' avant les bullets/tableau, et rédiger un contenu dense "
        "(15 à 40 sections pour un document professionnel complet)."
    ),
    parameters={
        "type": "object",
        "properties": {
            "document": {
                "type": "object",
                "description": (
                    "Structure du document (même format que export_docx). "
                    "Champs : title (str), sections (liste d'objets avec "
                    "heading, level, paragraphs, content, intro, bullets, table, page_break). "
                    "Rédiger un contenu complet et développé."
                )
            },
            "output_path": {
                "type": "string",
                "description": "Chemin de destination (ex: ~/Documents/rapport.pdf). Optionnel."
            },
            "filename": {
                "type": "string",
                "description": "Nom du fichier si output_path est omis."
            }
        },
        "required": ["document"]
    }
)
def export_pdf(document: dict, output_path: str = "", filename: str = "") -> str:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib import colors
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.lib.enums import TA_JUSTIFY, TA_LEFT
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
            PageBreak, HRFlowable, KeepTogether
        )

        name = filename or (document.get("title", "export") + ".pdf")
        if not name.endswith(".pdf"):
            name += ".pdf"
        p = _resolve_output(output_path, name)

        # ── Styles ────────────────────────────────────────────────────────
        BLUE = colors.HexColor("#2255A4")
        GREY = colors.HexColor("#F4F6FA")
        GLINE= colors.HexColor("#C5CDE0")

        def ST(n, **kw): return ParagraphStyle(n, **kw)
        styles = {
            "h0": ST("h0", fontName="Helvetica-Bold", fontSize=18, textColor=BLUE,
                     spaceBefore=0, spaceAfter=10),
            "h1": ST("h1", fontName="Helvetica-Bold", fontSize=14, textColor=BLUE,
                     spaceBefore=18, spaceAfter=6),
            "h2": ST("h2", fontName="Helvetica-Bold", fontSize=12,
                     textColor=colors.HexColor("#3B7ACC"), spaceBefore=12, spaceAfter=4),
            "h3": ST("h3", fontName="Helvetica-Bold", fontSize=11,
                     textColor=colors.HexColor("#374151"), spaceBefore=8, spaceAfter=3),
            "body": ST("body", fontName="Helvetica", fontSize=10,
                       leading=15, spaceBefore=2, spaceAfter=6, alignment=TA_JUSTIFY),
            "blt": ST("blt", fontName="Helvetica", fontSize=10, leading=14,
                      spaceBefore=2, spaceAfter=2, leftIndent=14),
            "th": ST("th", fontName="Helvetica-Bold", fontSize=9,
                     textColor=colors.white, leading=11),
            "tc": ST("tc", fontName="Helvetica", fontSize=9,
                     textColor=colors.HexColor("#111827"), leading=12),
        }

        # ── Pied de page ──────────────────────────────────────────────────
        def footer(canvas, doc):
            canvas.saveState()
            canvas.setFont("Helvetica", 7.5)
            canvas.setFillColor(colors.HexColor("#6B7280"))
            canvas.drawString(1.5*cm, 1*cm,
                              document.get("title", "Document généré par Prométhée AI"))
            canvas.drawRightString(A4[0] - 1.5*cm, 1*cm, f"Page {doc.page}")
            canvas.restoreState()

        # ── Assemblage ────────────────────────────────────────────────────
        story = []
        W = A4[0] - 3*cm  # largeur utile

        if document.get("title"):
            story.append(Paragraph(document["title"], styles["h0"]))
            story.append(HRFlowable(width="100%", thickness=2,
                                    color=BLUE, spaceAfter=8))

        for sec in document.get("sections", []):
            heading = sec.get("heading")
            level   = max(0, min(3, int(sec.get("level", 1))))
            if heading:
                style_key = f"h{level}" if level else "h1"
                story.append(Paragraph(heading, styles.get(style_key, styles["h1"])))

            # ── Contenu textuel (même logique que _build_docx) ────────────
            if sec.get("paragraphs"):
                for para in sec["paragraphs"]:
                    if para and str(para).strip():
                        story.append(Paragraph(str(para), styles["body"]))
            elif sec.get("content"):
                raw = sec["content"]
                if "\n\n" in raw:
                    parts = [p.strip() for p in raw.split("\n\n") if p.strip()]
                else:
                    parts = [p.strip() for p in raw.split("\n") if p.strip()]
                for part in parts:
                    story.append(Paragraph(part, styles["body"]))

            # ── Intro (avant bullets / tableau) ───────────────────────────
            if sec.get("intro"):
                story.append(Paragraph(str(sec["intro"]), styles["body"]))

            # ── Liste à puces ─────────────────────────────────────────────
            if sec.get("bullets"):
                for b in sec["bullets"]:
                    if b and str(b).strip():
                        story.append(Paragraph(
                            f'<font color="#2255A4">&#9658;</font>  {b}',
                            styles["blt"]
                        ))
                story.append(Spacer(1, 4))

            # ── Tableau ───────────────────────────────────────────────────
            if sec.get("table"):
                tbl_data = sec["table"]
                headers  = tbl_data.get("headers", [])
                rows     = tbl_data.get("rows", [])
                if headers:
                    ncols  = len(headers)
                    col_w  = W / ncols
                    data   = [[Paragraph(h, styles["th"]) for h in headers]]
                    for ri, row in enumerate(rows):
                        data.append([
                            Paragraph(str(v), styles["tc"])
                            for v in row[:ncols]
                        ])
                    ts = TableStyle([
                        ("BACKGROUND",   (0, 0), (-1,  0), BLUE),
                        ("ROWBACKGROUNDS",(0, 1), (-1, -1), [colors.white, GREY]),
                        ("GRID",         (0, 0), (-1, -1), 0.4, GLINE),
                        ("VALIGN",       (0, 0), (-1, -1), "TOP"),
                        ("TOPPADDING",   (0, 0), (-1, -1), 5),
                        ("BOTTOMPADDING",(0, 0), (-1, -1), 5),
                        ("LEFTPADDING",  (0, 0), (-1, -1), 6),
                        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                    ])
                    t = Table(data, colWidths=[col_w]*ncols, repeatRows=1, style=ts)
                    story.append(t)
                    story.append(Spacer(1, 8))

            if sec.get("page_break"):
                story.append(PageBreak())

        doc_obj = SimpleDocTemplate(
            str(p), pagesize=A4,
            topMargin=2*cm, bottomMargin=2.2*cm,
            leftMargin=1.5*cm, rightMargin=1.5*cm,
            title=document.get("title", ""),
        )
        doc_obj.build(story, onFirstPage=footer, onLaterPages=footer)

        return _ok(p, {"sections": len(document.get("sections", []))})
    except Exception as e:
        return _err(f"export_pdf : {e}")


# ═════════════════════════════════════════════════════════════════════════════
# LIBREOFFICE — CONVERSION ET NATIF
# ═════════════════════════════════════════════════════════════════════════════

def _libreoffice_convert(input_path: Path, target_format: str,
                          output_dir: Path) -> Path | None:
    """
    Appelle LibreOffice headless pour convertir un fichier.
    Retourne le chemin du fichier converti ou None en cas d'erreur.
    """
    soffice = shutil.which("libreoffice") or shutil.which("soffice")
    if not soffice:
        return None

    cmd = [
        soffice, "--headless", "--norestore",
        f"--convert-to", target_format,
        "--outdir", str(output_dir),
        str(input_path),
    ]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice : {result.stderr.strip()}")

    # LibreOffice place le fichier converti dans output_dir
    # avec le même stem que l'entrée
    ext_map = {
        "odt": ".odt", "ods": ".ods", "odp": ".odp",
        "pdf": ".pdf", "docx": ".docx", "xlsx": ".xlsx",
    }
    stem = input_path.stem
    ext  = ext_map.get(target_format, f".{target_format}")
    converted = output_dir / (stem + ext)
    if converted.exists():
        return converted

    # Chercher tout fichier avec la bonne extension (LibreOffice peut varier)
    candidates = list(output_dir.glob(f"{stem}*{ext}"))
    return candidates[0] if candidates else None


@tool(
    name="export_libreoffice",
    description=(
        "Convertit un fichier existant (docx, xlsx, pptx, pdf, csv…) vers un format "
        "LibreOffice natif : .odt (traitement de texte), .ods (tableur), .odp (présentation). "
        "Utilise LibreOffice en mode headless. Le fichier source doit exister sur le disque."
    ),
    parameters={
        "type": "object",
        "properties": {
            "input_path": {
                "type": "string",
                "description": "Chemin absolu ou relatif au home du fichier à convertir."
            },
            "target_format": {
                "type": "string",
                "enum": ["odt", "ods", "odp"],
                "description": "Format LibreOffice cible : odt (texte), ods (tableur), odp (présentation)."
            },
            "output_path": {
                "type": "string",
                "description": "Chemin du fichier de sortie. Si omis, placé à côté du fichier source."
            }
        },
        "required": ["input_path", "target_format"]
    }
)
def export_libreoffice(input_path: str, target_format: str,
                        output_path: str = "") -> str:
    try:
        src = Path(input_path).expanduser().resolve()
        if not src.exists():
            return _err(f"export_libreoffice : fichier source introuvable : {src}")

        # Dossier de sortie temporaire pour LibreOffice
        with tempfile.TemporaryDirectory() as tmp_dir:
            converted = _libreoffice_convert(src, target_format, Path(tmp_dir))
            if converted is None:
                return _err("export_libreoffice : LibreOffice non disponible ou conversion échouée")

            # Déterminer la destination finale
            if output_path:
                dest = Path(output_path).expanduser()
                if not dest.is_absolute():
                    dest = Path.home() / dest
            else:
                dest = src.parent / (src.stem + f".{target_format}")

            dest.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(str(converted), str(dest))

        return _ok(dest)
    except Exception as e:
        return _err(f"export_libreoffice : {e}")


@tool(
    name="export_libreoffice_native",
    description=(
        "Génère directement un fichier LibreOffice natif (.odt, .ods ou .odp) "
        "depuis une description JSON, sans passer par Word/Excel/PowerPoint. "
        "Le document est d'abord construit dans le format intermédiaire le plus "
        "approprié (docx→odt, xlsx→ods, pptx→odp) puis converti via LibreOffice. "
        "Même format JSON que export_docx / export_xlsx_json / export_pptx_json."
    ),
    parameters={
        "type": "object",
        "properties": {
            "target_format": {
                "type": "string",
                "enum": ["odt", "ods", "odp"],
                "description": "Format de sortie LibreOffice : odt (texte), ods (tableur), odp (présentation)."
            },
            "document": {
                "type": "object",
                "description": (
                    "Structure du document. Pour odt : même format que export_docx. "
                    "Pour ods : même format que export_xlsx_json (champ 'sheets'). "
                    "Pour odp : même format que export_pptx_json (champ 'slides')."
                )
            },
            "output_path": {
                "type": "string",
                "description": "Chemin de destination du fichier LibreOffice final. Optionnel."
            },
            "filename": {
                "type": "string",
                "description": "Nom du fichier si output_path est omis."
            }
        },
        "required": ["target_format", "document"]
    }
)
def export_libreoffice_native(target_format: str, document: dict,
                               output_path: str = "", filename: str = "") -> str:
    try:
        doc_title = document.get("title", "export")
        name = filename or f"{doc_title}.{target_format}"
        if not name.endswith(f".{target_format}"):
            name = name.rsplit(".", 1)[0] + f".{target_format}"
        final_dest = _resolve_output(output_path, name)

        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp = Path(tmp_dir)

            # ── Étape 1 : générer le format intermédiaire ──────────────────
            if target_format == "odt":
                tmp_src = tmp / f"{doc_title}.docx"
                doc = _build_docx(document)
                doc.save(str(tmp_src))
                intermediate_format = "odt"

            elif target_format == "ods":
                import openpyxl
                from openpyxl.styles import Font, PatternFill, Alignment
                from openpyxl.utils import get_column_letter

                tmp_src = tmp / f"{doc_title}.xlsx"
                wb = openpyxl.Workbook()
                wb.remove(wb.active)
                hf = Font(bold=True, color="FFFFFF")
                hfill = PatternFill("solid", fgColor="2255A4")
                ha = Alignment(horizontal="center")

                for sheet_def in document.get("sheets", []):
                    ws = wb.create_sheet(str(sheet_def.get("name", "Feuille"))[:31])
                    headers = sheet_def.get("headers", [])
                    rows    = sheet_def.get("rows", [])
                    for ci, h in enumerate(headers, 1):
                        c = ws.cell(row=1, column=ci, value=str(h))
                        c.font = hf; c.fill = hfill; c.alignment = ha
                    for ri, row in enumerate(rows, 2):
                        for ci, val in enumerate(row, 1):
                            ws.cell(row=ri, column=ci, value=val)
                    for col_cells in ws.columns:
                        max_len = max((len(str(c.value or "")) for c in col_cells), default=8)
                        ws.column_dimensions[get_column_letter(col_cells[0].column)].width = min(max_len + 4, 60)
                    ws.freeze_panes = "A2"
                wb.save(str(tmp_src))
                intermediate_format = "ods"

            elif target_format == "odp":
                from pptx import Presentation
                from pptx.util import Inches
                tmp_src = tmp / f"{doc_title}.pptx"
                prs = Presentation()
                prs.slide_width  = Inches(13.33)
                prs.slide_height = Inches(7.5)
                main_title = document.get("title", "")
                subtitle   = document.get("subtitle", "")
                if main_title:
                    _add_slide(prs, 0, main_title, [], subtitle=subtitle)
                for s in document.get("slides", []):
                    _add_slide(prs, 1, s.get("title", ""), s.get("bullets", []),
                               content=s.get("content", ""), notes=s.get("notes", ""))
                prs.save(str(tmp_src))
                intermediate_format = "odp"
            else:
                return _err(f"export_libreoffice_native : format non supporté : {target_format}")

            # ── Étape 2 : conversion via LibreOffice ───────────────────────
            converted = _libreoffice_convert(tmp_src, intermediate_format, tmp)
            if converted is None:
                return _err(
                    "export_libreoffice_native : LibreOffice non disponible. "
                    f"Le fichier intermédiaire est disponible : {tmp_src}"
                )

            final_dest.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(str(converted), str(final_dest))

        return _ok(final_dest, {"intermediate_format": intermediate_format})
    except Exception as e:
        return _err(f"export_libreoffice_native : {e}")
