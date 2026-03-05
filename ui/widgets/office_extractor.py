# ============================================================================
# Prométhée — Assistant IA desktop
# ============================================================================
# Auteur  : Pierre COUGET
# Licence : GNU Affero General Public License v3.0 (AGPL-3.0)
#           https://www.gnu.org/licenses/agpl-3.0.html
# Année   : 2026
# ----------------------------------------------------------------------------
# Ce fichier fait partie du projet Prométhée.
# Vous pouvez le redistribuer et/ou le modifier selon les termes de la
# licence AGPL-3.0 publiée par la Free Software Foundation.
# ============================================================================

"""
office_extractor.py — Extraction de contenu depuis fichiers Office (docx, xlsx, pptx)
"""
from pathlib import Path
from typing import Optional, Dict, Tuple


def is_available() -> bool:
    """Vérifie si les dépendances Office sont disponibles."""
    try:
        import docx
        import openpyxl
        from pptx import Presentation
        return True
    except ImportError:
        return False


def extract_from_docx(file_path: str) -> Tuple[Optional[str], Optional[str]]:
    """
    Extrait le texte d'un fichier Word (.docx).
    
    Args:
        file_path: Chemin vers le fichier .docx
    
    Returns:
        (contenu, erreur)
    """
    try:
        import docx
        
        doc = docx.Document(file_path)
        
        # Extraire le texte de tous les paragraphes
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)
        
        # Extraire le texte des tableaux
        tables_content = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_data.append(cell_text)
                if row_data:
                    table_data.append(" | ".join(row_data))
            if table_data:
                tables_content.append("\n".join(table_data))
        
        # Combiner tout
        content_parts = []
        
        if paragraphs:
            content_parts.append("\n\n".join(paragraphs))
        
        if tables_content:
            content_parts.append("\n\n=== TABLEAUX ===\n\n" + "\n\n".join(tables_content))
        
        if not content_parts:
            return None, "Document Word vide"
        
        content = "\n\n".join(content_parts)
        
        # Limiter la taille
        if len(content) > 100000:
            content = content[:100000] + "\n\n... [contenu tronqué]"
        
        return content, None
    
    except ImportError:
        return None, "python-docx non installé (pip install python-docx)"
    
    except Exception as e:
        return None, f"Erreur lecture Word : {str(e)}"


def extract_from_xlsx(file_path: str, max_sheets: int = 5) -> Tuple[Optional[str], Optional[str]]:
    """
    Extrait le contenu d'un fichier Excel (.xlsx).
    
    Args:
        file_path: Chemin vers le fichier .xlsx
        max_sheets: Nombre maximum de feuilles à extraire
    
    Returns:
        (contenu, erreur)
    """
    try:
        import openpyxl
        
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        
        sheets_content = []
        sheet_names = workbook.sheetnames[:max_sheets]
        
        for sheet_name in sheet_names:
            sheet = workbook[sheet_name]
            
            # Extraire les données
            rows_data = []
            for row in sheet.iter_rows(values_only=True):
                # Filtrer les lignes vides
                row_values = [str(cell) if cell is not None else "" for cell in row]
                if any(val.strip() for val in row_values):
                    rows_data.append(" | ".join(row_values))
            
            if rows_data:
                sheet_content = f"=== Feuille: {sheet_name} ===\n\n"
                sheet_content += "\n".join(rows_data)
                sheets_content.append(sheet_content)
        
        workbook.close()
        
        if not sheets_content:
            return None, "Classeur Excel vide"
        
        content = "\n\n".join(sheets_content)
        
        # Ajouter un résumé
        total_sheets = len(workbook.sheetnames)
        if total_sheets > max_sheets:
            content = f"[Affichage de {max_sheets}/{total_sheets} feuilles]\n\n" + content
        
        # Limiter la taille
        if len(content) > 100000:
            content = content[:100000] + "\n\n... [contenu tronqué]"
        
        return content, None
    
    except ImportError:
        return None, "openpyxl non installé (pip install openpyxl)"
    
    except Exception as e:
        return None, f"Erreur lecture Excel : {str(e)}"


def extract_from_pptx(file_path: str, max_slides: int = 20) -> Tuple[Optional[str], Optional[str]]:
    """
    Extrait le texte d'une présentation PowerPoint (.pptx).
    
    Args:
        file_path: Chemin vers le fichier .pptx
        max_slides: Nombre maximum de slides à extraire
    
    Returns:
        (contenu, erreur)
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        
        slides_content = []
        total_slides = len(prs.slides)
        slides_to_process = min(total_slides, max_slides)
        
        for i, slide in enumerate(prs.slides[:slides_to_process], 1):
            slide_text = []
            
            # Extraire le texte de toutes les formes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # Extraire le texte des tableaux
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_data.append(cell.text.strip())
                        if row_data:
                            slide_text.append(" | ".join(row_data))
            
            if slide_text:
                slide_content = f"=== Slide {i} ===\n\n"
                slide_content += "\n\n".join(slide_text)
                slides_content.append(slide_content)
        
        if not slides_content:
            return None, "Présentation PowerPoint vide"
        
        content = "\n\n".join(slides_content)
        
        # Ajouter un résumé
        if total_slides > max_slides:
            content = f"[Affichage de {max_slides}/{total_slides} slides]\n\n" + content
        
        # Limiter la taille
        if len(content) > 100000:
            content = content[:100000] + "\n\n... [contenu tronqué]"
        
        return content, None
    
    except ImportError:
        return None, "python-pptx non installé (pip install python-pptx)"
    
    except Exception as e:
        return None, f"Erreur lecture PowerPoint : {str(e)}"


def extract_office_metadata(file_path: str) -> Dict[str, any]:
    """
    Extrait les métadonnées d'un fichier Office.
    
    Args:
        file_path: Chemin vers le fichier Office
    
    Returns:
        Dict avec les métadonnées disponibles
    """
    metadata = {
        "format": Path(file_path).suffix.lower(),
        "size": Path(file_path).stat().st_size,
    }
    
    try:
        suffix = Path(file_path).suffix.lower()
        
        if suffix == ".docx":
            import docx
            doc = docx.Document(file_path)
            core_props = doc.core_properties
            
            metadata.update({
                "author": core_props.author or None,
                "title": core_props.title or None,
                "subject": core_props.subject or None,
                "created": core_props.created.isoformat() if core_props.created else None,
                "modified": core_props.modified.isoformat() if core_props.modified else None,
                "paragraphs": len(doc.paragraphs),
                "tables": len(doc.tables),
            })
        
        elif suffix == ".xlsx":
            import openpyxl
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            
            metadata.update({
                "sheets": len(workbook.sheetnames),
                "sheet_names": workbook.sheetnames[:5],  # Premières 5 feuilles
            })
            workbook.close()
        
        elif suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            
            metadata.update({
                "slides": len(prs.slides),
            })
            
            if hasattr(prs.core_properties, 'author'):
                metadata["author"] = prs.core_properties.author
            if hasattr(prs.core_properties, 'title'):
                metadata["title"] = prs.core_properties.title
    
    except Exception:
        pass  # Les métadonnées sont optionnelles
    
    return metadata


def get_office_summary(file_path: str) -> str:
    """
    Génère un résumé formaté d'un fichier Office.
    
    Args:
        file_path: Chemin vers le fichier Office
    
    Returns:
        Résumé formaté
    """
    suffix = Path(file_path).suffix.lower()
    file_name = Path(file_path).name
    
    # Extraire le contenu
    content = None
    error = None
    
    if suffix == ".docx":
        content, error = extract_from_docx(file_path)
    elif suffix == ".xlsx":
        content, error = extract_from_xlsx(file_path)
    elif suffix == ".pptx":
        content, error = extract_from_pptx(file_path)
    else:
        return f"Format non supporté : {suffix}"
    
    if error:
        return f"Erreur : {error}"
    
    # Extraire les métadonnées
    metadata = extract_office_metadata(file_path)
    
    # Formater le résumé
    summary_parts = [f"# {file_name}\n"]
    
    # Ajouter les métadonnées pertinentes
    meta_info = []
    if metadata.get("author"):
        meta_info.append(f"**Auteur** : {metadata['author']}")
    if metadata.get("title") and metadata["title"] != file_name:
        meta_info.append(f"**Titre** : {metadata['title']}")
    if metadata.get("paragraphs"):
        meta_info.append(f"**Paragraphes** : {metadata['paragraphs']}")
    if metadata.get("tables"):
        meta_info.append(f"**Tableaux** : {metadata['tables']}")
    if metadata.get("sheets"):
        meta_info.append(f"**Feuilles** : {metadata['sheets']}")
    if metadata.get("slides"):
        meta_info.append(f"**Slides** : {metadata['slides']}")
    
    if meta_info:
        summary_parts.append(" | ".join(meta_info) + "\n")
    
    # Ajouter le contenu
    if content:
        summary_parts.append(f"\n{content}")
    
    return "\n".join(summary_parts)
